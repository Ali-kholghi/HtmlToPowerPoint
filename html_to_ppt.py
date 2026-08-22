import ctypes
import io
import os
import pathlib
import sys
import threading
import time
import tkinter as tk
from tkinter import filedialog, messagebox

# Enable High-DPI Awareness on Windows to eliminate GUI blurriness
try:
    if sys.platform == "win32":
        ctypes.windll.shcore.SetProcessDpiAwareness(2)  # Per-monitor DPI aware
except Exception:
    try:
        ctypes.windll.user32.SetProcessDPIAware()
    except Exception:
        pass

import customtkinter as ctk
from PIL import Image
from playwright.sync_api import sync_playwright
from pptx import Presentation
from pptx.util import Inches

ctk.set_appearance_mode("Dark")
ctk.set_default_color_theme("blue")


class UniversalHTMLConverter:
    """
    Production-grade converter that preserves native HTML layouts,
    prevents SVG/grid distortion, and mathematically centers every slide.
    """

    def __init__(self, status_callback, log_callback, progress_callback):
        self.status_cb = status_callback
        self.log_cb = log_callback
        self.progress_cb = progress_callback
        self.is_cancelled = False

    def cancel(self):
        self.is_cancelled = True

    def convert(self, html_path, ppt_path, mode="auto", scale_factor=2, aspect_ratio="16:9", compression="opt_png"):
        self.is_cancelled = False
        self.log_cb(f"Starting conversion for: {os.path.basename(html_path)}")
        self.status_cb("Initializing Chromium engine...")
        self.progress_cb(0.05)

        viewport_w = 1920
        viewport_h = 1080 if aspect_ratio == "16:9" else 1440

        with sync_playwright() as p:
            browser = p.chromium.launch(
                headless=True,
                args=[
                    "--disable-web-security",
                    "--allow-file-access-from-files",
                ],
            )

            context = browser.new_context(
                viewport={"width": viewport_w, "height": viewport_h},
                device_scale_factor=scale_factor,
            )
            page = context.new_page()

            # 1. Load HTML
            self.status_cb("Loading document...")
            file_url = pathlib.Path(html_path).absolute().as_uri()
            page.goto(file_url, wait_until="networkidle", timeout=60000)

            # 2. Universal Pre-Flight
            self.status_cb("Stabilizing web fonts & styles...")
            self._universal_preflight(page)
            self.progress_cb(0.15)

            # 3. Strategy Resolution
            selected_strategy = mode
            if mode == "auto":
                selected_strategy = self._detect_strategy(page)
                self.log_cb(f"Auto-detected strategy: [{selected_strategy.upper()}]")
            else:
                self.log_cb(f"Selected strategy: [{selected_strategy.upper()}]")

            # 4. Smart Capture
            raw_screenshots = []
            if selected_strategy in ["presentation_deck", "section"]:
                raw_screenshots = self._capture_element_slides(page)
            elif selected_strategy == "revealjs":
                raw_screenshots = self._capture_revealjs(page)
            elif selected_strategy == "marp":
                raw_screenshots = self._capture_marp(page)
            else:
                raw_screenshots = self._capture_smart_dom_scroll(page, viewport_h)

            browser.close()

            if self.is_cancelled:
                self.log_cb("Conversion cancelled by user.")
                self.status_cb("Cancelled")
                return

            if not raw_screenshots:
                raise RuntimeError("No visual content could be captured.")

            # 5. Image Optimization
            self.status_cb("Optimizing image quality & file size...")
            self.progress_cb(0.85)
            processed_images = self._optimize_images(raw_screenshots, compression)

            # 6. Assemble PowerPoint with Precision Centering
            self.status_cb("Assembling PowerPoint slides...")
            self.progress_cb(0.92)
            self._build_powerpoint(processed_images, ppt_path, aspect_ratio)

            self.progress_cb(1.0)
            self.status_cb("Conversion complete!")
            self.log_cb(f"Saved: {ppt_path}")

    def _universal_preflight(self, page):
        """Universal stabilization: disables transitions, scrollbars, and waits for web fonts."""
        page.evaluate("""() => {
            const style = document.createElement('style');
            style.id = '__universal_converter_styles';
            style.textContent = `
                ::-webkit-scrollbar { display: none !important; }
                * { 
                    scrollbar-width: none !important; 
                    -ms-overflow-style: none !important;
                    animation-duration: 0s !important;
                    animation-delay: 0s !important;
                    transition-duration: 0s !important;
                    transition-delay: 0s !important;
                }
            `;
            document.head.appendChild(style);
        }""")

        try:
            page.evaluate("() => document.fonts.ready")
        except Exception:
            pass

        # Trigger lazy-loaded assets
        page.evaluate("() => { window.scrollTo(0, document.body.scrollHeight); }")
        time.sleep(0.2)
        page.evaluate("() => { window.scrollTo(0, 0); }")
        time.sleep(0.1)

    def _detect_strategy(self, page):
        return page.evaluate("""() => {
            if (window.Reveal && typeof window.Reveal.isReady === 'function') return 'revealjs';
            if (document.querySelector('.marpit') || document.querySelector('.marp-slide')) return 'marp';
            if (document.querySelector('.slide-section, [data-title], .slide') && typeof window.togglePresentationMode === 'function') {
                return 'presentation_deck';
            }
            const sections = document.querySelectorAll('section, article, .page, .slide, .slide-section');
            if (sections.length > 1) return 'section';
            return 'scroll';
        }""")

    def _capture_element_slides(self, page):
        """
        Captures each slide element directly without altering its natural design,
        preventing SVG blowups, squished multi-column grids, and clipping.
        """
        self.log_cb("Using Element-Centric Capture Frame...")

        # Hide floating toolbars/navbars that might contaminate captures
        page.evaluate("""() => {
            if (typeof togglePresentationMode === 'function' && !document.body.classList.contains('presentation-mode')) {
                togglePresentationMode();
            }
            const nav = document.querySelector('.presenter-nav, header, nav');
            if (nav) nav.style.display = 'none';

            // Remove artificial min-height locks that force vertical voids
            document.querySelectorAll('.slide-section, .slide, section').forEach(el => {
                el.style.minHeight = 'auto';
                el.style.marginBottom = '0px';
            });
        }""")
        time.sleep(0.2)

        slides = page.locator(".slide-section, .slide, section, article")
        count = slides.count()
        screenshots = []

        for i in range(count):
            if self.is_cancelled:
                break
            self.status_cb(f"Capturing slide {i+1} of {count}...")
            self.progress_cb(0.2 + (0.6 * (i / max(1, count))))

            # Make only current slide visible and scroll into view
            page.evaluate(f"""(idx) => {{
                if (typeof showSlide === 'function') {{
                    showSlide(idx);
                }} else {{
                    const els = document.querySelectorAll('.slide-section, .slide, section, article');
                    els.forEach((el, j) => {{
                        el.style.display = (j === idx) ? 'block' : 'none';
                    }});
                }}
                window.scrollTo(0, 0);
            }}""", i)
            time.sleep(0.2)

            target_slide = page.locator(".active-slide").first
            if target_slide.count() == 0:
                target_slide = slides.nth(i)

            # Capture the element directly (includes borders, shadows, backgrounds)
            target_slide.scroll_into_view_if_needed()
            screenshots.append(target_slide.screenshot(type="png"))
            self.log_cb(f"Slide {i+1}: Captured in natural design layout.")

        return screenshots

    def _capture_revealjs(self, page):
        total_slides = page.evaluate("() => Reveal.getTotalSlides()")
        screenshots = []
        for i in range(total_slides):
            if self.is_cancelled:
                break
            self.status_cb(f"Capturing slide {i+1} of {total_slides}...")
            self.progress_cb(0.2 + (0.6 * (i / total_slides)))
            page.evaluate(f"() => Reveal.slide({i})")
            time.sleep(0.25)
            screenshots.append(page.screenshot(type="png", full_page=False))
        return screenshots

    def _capture_marp(self, page):
        slides = page.locator(".marpit > svg, .marp-slide")
        count = slides.count()
        screenshots = []
        for i in range(count):
            if self.is_cancelled:
                break
            self.status_cb(f"Capturing Marp slide {i+1} of {count}...")
            self.progress_cb(0.2 + (0.6 * (i / count)))
            slide = slides.nth(i)
            slide.scroll_into_view_if_needed()
            time.sleep(0.2)
            screenshots.append(slide.screenshot(type="png"))
        return screenshots

    def _capture_smart_dom_scroll(self, page, viewport_h):
        """DOM-Aware Slicing: Only cuts between paragraphs/tables for continuous pages."""
        self.log_cb("Using DOM-Aware Slicing for continuous page...")
        total_height = page.evaluate("() => document.documentElement.scrollHeight || document.body.scrollHeight")
        current_scroll = 0
        page_idx = 1
        screenshots = []

        while current_scroll < total_height:
            if self.is_cancelled:
                break
            self.status_cb(f"Processing page slice {page_idx}...")

            if page_idx > 1:
                page.evaluate("""() => {
                    document.querySelectorAll('*').forEach(el => {
                        const style = window.getComputedStyle(el);
                        if (style.position === 'fixed' || style.position === 'sticky') {
                            el.style.visibility = 'hidden';
                        }
                    });
                }""")

            page.evaluate(f"() => window.scrollTo(0, {current_scroll})")
            time.sleep(0.2)

            screenshots.append(page.screenshot(type="png", full_page=False))

            next_scroll = page.evaluate(f"""(curr) => {{
                const targetY = curr + {viewport_h};
                const totalH = Math.max(document.body.scrollHeight, document.documentElement.scrollHeight);
                if (targetY >= totalH) return totalH;

                const elements = document.querySelectorAll('p, h1, h2, h3, h4, h5, h6, tr, li, .card, section, article, div, table');
                let bestCutY = targetY;
                const minAcceptableY = curr + ({viewport_h} * 0.70);

                for (let el of elements) {{
                    const rect = el.getBoundingClientRect();
                    const elTop = rect.top + window.scrollY;
                    const elBottom = rect.bottom + window.scrollY;

                    if (elTop < targetY && elBottom > targetY) {{
                        if (elTop > minAcceptableY) {{
                            bestCutY = Math.min(bestCutY, elTop - 10);
                        }}
                    }}
                }}
                return Math.floor(bestCutY);
            }}""", current_scroll)

            if next_scroll <= current_scroll:
                next_scroll = current_scroll + viewport_h

            current_scroll = next_scroll
            page_idx += 1

        return screenshots

    def _optimize_images(self, raw_screenshots, compression):
        processed = []
        for img_bytes in raw_screenshots:
            with Image.open(io.BytesIO(img_bytes)) as img:
                out = io.BytesIO()
                
                if compression == "opt_png":
                    bg = Image.new("RGB", img.size, (255, 255, 255))
                    if img.mode == "RGBA":
                        bg.paste(img, mask=img.split()[3])
                    else:
                        bg.paste(img)
                    quantized = bg.quantize(colors=256, method=Image.Quantize.MEDIANCUT)
                    quantized.save(out, format="PNG", optimize=True)
                
                elif compression == "jpeg_444":
                    bg = Image.new("RGB", img.size, (255, 255, 255))
                    if img.mode == "RGBA":
                        bg.paste(img, mask=img.split()[3])
                    else:
                        bg.paste(img)
                    bg.save(out, format="JPEG", quality=94, subsampling=0, optimize=True)
                
                else:  # Lossless Full PNG
                    img.save(out, format="PNG", optimize=True)

                processed.append(out.getvalue())
        return processed

    def _build_powerpoint(self, screenshots, ppt_path, aspect_ratio):
        """
        Builds the PowerPoint deck with mathematical 2D centering.
        Leaves aesthetic breathing margins so slides look professional.
        """
        prs = Presentation()

        if aspect_ratio == "16:9":
            slide_w_in, slide_h_in = 13.333, 7.5
        else:
            slide_w_in, slide_h_in = 10.0, 7.5

        prs.slide_width = Inches(slide_w_in)
        prs.slide_height = Inches(slide_h_in)
        blank_layout = prs.slide_layouts[6]

        # Use 94% of slide bounds for padding
        max_target_w = slide_w_in * 0.94
        max_target_h = slide_h_in * 0.94

        for img_bytes in screenshots:
            slide = prs.slides.add_slide(blank_layout)
            image_stream = io.BytesIO(img_bytes)

            with Image.open(image_stream) as img:
                img_w, img_h = img.size
                img_ratio = img_w / img_h

                # Fit within available box preserving aspect ratio
                if (max_target_w / img_ratio) <= max_target_h:
                    disp_w = max_target_w
                    disp_h = max_target_w / img_ratio
                else:
                    disp_h = max_target_h
                    disp_w = max_target_h * img_ratio

                # Mathematical centering on the PowerPoint canvas
                left = (slide_w_in - disp_w) / 2
                top = (slide_h_in - disp_h) / 2

            image_stream.seek(0)
            slide.shapes.add_picture(
                image_stream,
                Inches(left),
                Inches(top),
                width=Inches(disp_w),
                height=Inches(disp_h),
            )

        prs.save(ppt_path)


# --- Modern Desktop GUI ---

class ModernApp(ctk.CTk):
    def __init__(self):
        super().__init__()

        self.title("Universal HTML to PowerPoint Studio v5")
        self.geometry("820x660")
        self.minsize(720, 580)

        self.converter = None
        self.worker_thread = None

        self._build_ui()

    def _build_ui(self):
        header_frame = ctk.CTkFrame(self, corner_radius=12, fg_color=("gray85", "gray17"))
        header_frame.pack(fill="x", padx=20, pady=(15, 10))

        title_lbl = ctk.CTkLabel(
            header_frame,
            text="✨ Universal HTML → PowerPoint Studio",
            font=ctk.CTkFont(size=20, weight="bold"),
        )
        title_lbl.pack(anchor="w", padx=16, pady=(10, 2))

        subtitle_lbl = ctk.CTkLabel(
            header_frame,
            text="Element-Centric Framing • Responsive Grid & SVG Protection • Precision 2D Centering",
            font=ctk.CTkFont(size=12),
            text_color="gray",
        )
        subtitle_lbl.pack(anchor="w", padx=16, pady=(0, 10))

        card = ctk.CTkFrame(self, corner_radius=12)
        card.pack(fill="x", padx=20, pady=5)

        file_grid = ctk.CTkFrame(card, fg_color="transparent")
        file_grid.pack(fill="x", padx=15, pady=12)
        file_grid.columnconfigure(1, weight=1)

        ctk.CTkLabel(file_grid, text="Source HTML:", font=ctk.CTkFont(weight="bold")).grid(row=0, column=0, sticky="w", pady=5)
        self.html_entry = ctk.CTkEntry(file_grid, placeholder_text="Select or drop any HTML file...")
        self.html_entry.grid(row=0, column=1, sticky="ew", padx=10, pady=5)
        self.browse_html_btn = ctk.CTkButton(file_grid, text="Browse", width=90, command=self.browse_html)
        self.browse_html_btn.grid(row=0, column=2, pady=5)

        ctk.CTkLabel(file_grid, text="Output PPTX:", font=ctk.CTkFont(weight="bold")).grid(row=1, column=0, sticky="w", pady=5)
        self.ppt_entry = ctk.CTkEntry(file_grid, placeholder_text="Path for the generated presentation...")
        self.ppt_entry.grid(row=1, column=1, sticky="ew", padx=10, pady=5)
        self.browse_ppt_btn = ctk.CTkButton(file_grid, text="Browse", width=90, command=self.browse_ppt)
        self.browse_ppt_btn.grid(row=1, column=2, pady=5)

        # Options Row
        options_frame = ctk.CTkFrame(card, fg_color="transparent")
        options_frame.pack(fill="x", padx=15, pady=(0, 15))

        # 1. Strategy
        strat_box = ctk.CTkFrame(options_frame, fg_color=("gray90", "gray20"), corner_radius=8)
        strat_box.pack(side="left", fill="both", expand=True, padx=(0, 5), pady=4)
        ctk.CTkLabel(strat_box, text="Capture Strategy", font=ctk.CTkFont(size=11, weight="bold")).pack(anchor="w", padx=10, pady=(6, 2))
        self.strategy_var = ctk.StringVar(value="auto")
        self.strategy_menu = ctk.CTkOptionMenu(
            strat_box,
            values=["auto", "presentation_deck", "section", "scroll", "revealjs", "marp"],
            variable=self.strategy_var,
        )
        self.strategy_menu.pack(fill="x", padx=10, pady=(0, 8))

        # 2. Quality
        quality_box = ctk.CTkFrame(options_frame, fg_color=("gray90", "gray20"), corner_radius=8)
        quality_box.pack(side="left", fill="both", expand=True, padx=5, pady=4)
        ctk.CTkLabel(quality_box, text="Rendering Quality", font=ctk.CTkFont(size=11, weight="bold")).pack(anchor="w", padx=10, pady=(6, 2))
        self.scale_var = ctk.StringVar(value="2x (Ultra Crisp 4K)")
        self.scale_menu = ctk.CTkOptionMenu(
            quality_box,
            values=["1x (Standard 1080p)", "2x (Ultra Crisp 4K)", "3x (Print Grade)"],
            variable=self.scale_var,
        )
        self.scale_menu.pack(fill="x", padx=10, pady=(0, 8))

        # 3. Compression
        comp_box = ctk.CTkFrame(options_frame, fg_color=("gray90", "gray20"), corner_radius=8)
        comp_box.pack(side="left", fill="both", expand=True, padx=(5, 0), pady=4)
        ctk.CTkLabel(comp_box, text="File Size Optimizer", font=ctk.CTkFont(size=11, weight="bold")).pack(anchor="w", padx=10, pady=(6, 2))
        self.comp_var = ctk.StringVar(value="Optimized PNG (70% Smaller)")
        self.comp_menu = ctk.CTkOptionMenu(
            comp_box,
            values=["Optimized PNG (70% Smaller)", "High-Q JPEG 4:4:4 (85% Smaller)", "Lossless Full PNG"],
            variable=self.comp_var,
        )
        self.comp_menu.pack(fill="x", padx=10, pady=(0, 8))

        # Progress & Action Area
        action_frame = ctk.CTkFrame(self, fg_color="transparent")
        action_frame.pack(fill="x", padx=20, pady=5)

        self.progress_bar = ctk.CTkProgressBar(action_frame)
        self.progress_bar.pack(fill="x", pady=(0, 6))
        self.progress_bar.set(0)

        status_row = ctk.CTkFrame(action_frame, fg_color="transparent")
        status_row.pack(fill="x")

        self.status_label = ctk.CTkLabel(status_row, text="Status: Ready", font=ctk.CTkFont(size=12, slant="italic"))
        self.status_label.pack(side="left")

        self.convert_btn = ctk.CTkButton(
            status_row,
            text="Convert to PowerPoint",
            font=ctk.CTkFont(weight="bold"),
            command=self.start_conversion,
            height=36,
            width=180,
        )
        self.convert_btn.pack(side="right")

        # Live Console Output Box
        log_frame = ctk.CTkFrame(self, corner_radius=10)
        log_frame.pack(fill="both", expand=True, padx=20, pady=(8, 15))

        log_hdr = ctk.CTkLabel(log_frame, text="Activity Log", font=ctk.CTkFont(size=11, weight="bold"))
        log_hdr.pack(anchor="w", padx=12, pady=(4, 2))

        self.log_textbox = ctk.CTkTextbox(log_frame, font=ctk.CTkFont(family="Consolas", size=11))
        self.log_textbox.pack(fill="both", expand=True, padx=10, pady=(0, 8))
        self.log_textbox.configure(state="disabled")

    def browse_html(self):
        file_path = filedialog.askopenfilename(
            title="Select HTML File",
            filetypes=[("HTML Files", "*.html *.htm"), ("All Files", "*.*")],
        )
        if file_path:
            self.html_entry.delete(0, tk.END)
            self.html_entry.insert(0, file_path)
            suggested_ppt = os.path.splitext(file_path)[0] + ".pptx"
            self.ppt_entry.delete(0, tk.END)
            self.ppt_entry.insert(0, suggested_ppt)

    def browse_ppt(self):
        file_path = filedialog.asksaveasfilename(
            title="Save Presentation As",
            defaultextension=".pptx",
            filetypes=[("PowerPoint Presentations", "*.pptx")],
        )
        if file_path:
            self.ppt_entry.delete(0, tk.END)
            self.ppt_entry.insert(0, file_path)

    def log(self, text):
        self.log_textbox.configure(state="normal")
        self.log_textbox.insert("end", f"[{time.strftime('%H:%M:%S')}] {text}\n")
        self.log_textbox.see("end")
        self.log_textbox.configure(state="disabled")

    def update_status(self, text):
        self.status_label.configure(text=f"Status: {text}")

    def update_progress(self, val):
        self.progress_bar.set(val)

    def set_ui_state(self, running=True):
        state = "disabled" if running else "normal"
        self.convert_btn.configure(state=state)
        self.browse_html_btn.configure(state=state)
        self.browse_ppt_btn.configure(state=state)
        self.strategy_menu.configure(state=state)
        self.scale_menu.configure(state=state)
        self.comp_menu.configure(state=state)

    def start_conversion(self):
        html_file = self.html_entry.get().strip()
        ppt_file = self.ppt_entry.get().strip()
        strategy = self.strategy_var.get()

        scale_text = self.scale_var.get()
        scale = 2
        if "1x" in scale_text:
            scale = 1
        elif "3x" in scale_text:
            scale = 3

        comp_text = self.comp_var.get()
        compression = "opt_png"
        if "JPEG" in comp_text:
            compression = "jpeg_444"
        elif "Lossless" in comp_text:
            compression = "lossless_png"

        if not html_file or not os.path.exists(html_file):
            messagebox.showerror("Error", "Please select a valid existing HTML file.")
            return
        if not ppt_file:
            messagebox.showerror("Error", "Please specify an output location for the .pptx file.")
            return

        self.set_ui_state(running=True)
        self.update_progress(0)
        self.log_textbox.configure(state="normal")
        self.log_textbox.delete("1.0", "end")
        self.log_textbox.configure(state="disabled")

        self.converter = UniversalHTMLConverter(
            status_callback=lambda s: self.after(0, self.update_status, s),
            log_callback=lambda l: self.after(0, self.log, l),
            progress_callback=lambda p: self.after(0, self.update_progress, p),
        )

        def worker():
            try:
                self.converter.convert(
                    html_path=html_file,
                    ppt_path=ppt_file,
                    mode=strategy,
                    scale_factor=scale,
                    aspect_ratio="16:9",
                    compression=compression,
                )
                self.after(0, lambda: messagebox.showinfo("Success", f"Presentation generated successfully!\n\nLocation: {ppt_file}"))
            except Exception as e:
                self.after(0, lambda: messagebox.showerror("Conversion Failed", f"An error occurred:\n\n{str(e)}"))
                self.after(0, self.log, f"ERROR: {str(e)}")
            finally:
                self.after(0, self.set_ui_state, False)

        self.worker_thread = threading.Thread(target=worker, daemon=True)
        self.worker_thread.start()


if __name__ == "__main__":
    app = ModernApp()
    app.mainloop()